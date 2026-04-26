#!/usr/bin/env python3
"""
Human AI: ConverterAgent - v1.0
Specialized agent for transforming multi-format documents (PDF, DOCX, PPTX, JSON) 
into standardized Markdown/Text for the Researcher and Developer agents.
"""

import os
import json
import re
from typing import Dict, Any, Optional
from pathlib import Path

# Third-party library imports (installed in venv)
try:
    import PyPDF2
    import docx
    import pptx
    from markitdown import MarkItDown
except ImportError:
    # These will be installed via the venv during the first run if missing
    pass

class ConverterAgent:
    def __init__(self, output_dir: str = "/home/yahwehatwork/human-ai/docs/converted"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def convert(self, file_path: str) -> Optional[str]:
        """
        Main entry point for conversion.
        Returns the path to the converted markdown file.
        """
        if not os.path.exists(file_path):
            print(f"❌ Error: File {file_path} not found.")
            return None

        ext = Path(file_path).suffix.lower()
        filename = Path(file_path).stem
        target_path = os.path.join(self.output_dir, f"{filename}_converted.md")

        try:
            if ext == '.pdf':
                content = self._convert_pdf(file_path)
            elif ext == '.docx':
                content = self._convert_docx(file_path)
            elif ext == '.pptx':
                content = self._convert_pptx(file_path)
            elif ext == '.json':
                content = self._convert_json(file_path)
            elif ext == '.txt' or ext == '.md':
                content = self._read_text(file_path)
            else:
                print(f"⚠️ Unsupported format {ext}. Attempting raw text read.")
                content = self._read_text(file_path)

            if content:
                with open(target_path, 'w', encoding='utf-8') as f:
                    f.write(content)
                return target_path
            
        except Exception as e:
            print(f"❌ Conversion error for {file_path}: {e}")
        
        return None

    def _read_text(self, path: str) -> str:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def _convert_pdf(self, path: str) -> str:
        import PyPDF2
        text = ""
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text

    def _convert_docx(self, path: str) -> str:
        import docx
        doc = docx.Document(path)
        return "\n".join([para.text for para in doc.paragraphs])

    def _convert_pptx(self, path: str) -> str:
        import pptx
        prs = pptx.Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def _convert_json(self, path: str) -> str:
        with open(path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return json.dumps(data, indent=2)

    async def close(self):
        """Placeholder for async cleanup if needed."""
        pass

if __name__ == "__main__":
    # Basic internal test
    converter = ConverterAgent()
    print("ConverterAgent initialized. Ready to process documents.")
